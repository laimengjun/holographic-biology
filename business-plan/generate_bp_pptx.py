# -*- coding: utf-8 -*-
"""
HoloScan 2.0 投资人 BP · PPTX 生成器

基于 python-pptx，生成 20 页投资人路演 PPT
设计：深色背景 + 红色强调（与 HoloScan 品牌一致）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

# 颜色（品牌色）
BRAND_RED = RGBColor(0xC9, 0x2A, 0x2A)      # 主红色
DARK_BG = RGBColor(0x1A, 0x1A, 0x1A)         # 深背景
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)       # 浅背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # 白色
GRAY = RGBColor(0x6C, 0x75, 0x7D)           # 灰色
ACCENT_YELLOW = RGBColor(0xFF, 0xD8, 0x43)  # 强调黄
ACCENT_BLUE = RGBColor(0x4D, 0xAB, 0xF7)    # 强调蓝

# 输出路径
OUTPUT = r"D:\obsidian\Holographic-Biology\business-plan\holoscan-2.0-investor-bp.pptx"

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # ===== Slide 1: 封面 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_bg(slide, DARK_BG)

    # 大标题
    add_text(slide, "HoloScan 2.0", Inches(1), Inches(2), Inches(11), Inches(1.5),
             font_size=72, bold=True, color=BRAND_RED, align=PP_ALIGN.CENTER)
    add_text(slide, "让每一次体检都更早一步", Inches(1), Inches(3.5), Inches(11), Inches(0.8),
             font_size=32, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "基于张颖清全息胚理论的 8 微系统 AI 健康监测平台", Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

    # 底部
    add_text(slide, "2026 Q3 · 种子轮融资 · ¥500 万", Inches(1), Inches(6.3), Inches(11), Inches(0.4),
             font_size=14, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)
    add_text(slide, "厦门 | 2026", Inches(1), Inches(6.8), Inches(11), Inches(0.3),
             font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # ===== Slide 2: 问题 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_text(slide, "问题：三大痛点", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    # 3 个痛点卡片
    pain_points = [
        ("癌症早期检出率仅 30-40%", "中国每年新增癌症 457 万例，60% 发现时已是中晚期"),
        ("中医缺乏客观化工具", "3000 年临床经验难以规模化推广"),
        ("缺乏多模态融合监测", "单一传感器数据有限"),
    ]
    for i, (title, desc) in enumerate(pain_points):
        x = Inches(0.5 + i * 4.3)
        add_box(slide, x, Inches(1.5), Inches(4), Inches(4.5),
                fill=BRAND_RED, line=None)
        add_text(slide, f"#{i+1}", x, Inches(1.8), Inches(4), Inches(0.6),
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, x, Inches(2.5), Inches(4), Inches(1.5),
                 font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x, Inches(4.0), Inches(4), Inches(1.5),
                 font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    # 底部数据
    add_text(slide, "中国癌症 5 年生存率 40.5% vs 美国 67% | 慢病管理 ¥3000 亿 | 健康监测硬件 ¥5000 亿",
             Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # ===== Slide 3: 解决方案 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_text(slide, "HoloScan 2.0", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)
    add_text(slide, "8 微系统 + AI 融合 + 数字孪生", Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    # 流程图
    flow_steps = ["用户", "8 微系统\n采集", "多模态融合\nTransformer", "数字孪生\n引擎", "中医辨证\nAI", "风险\n预警"]
    for i, step in enumerate(flow_steps):
        x = Inches(0.4 + i * 2.1)
        add_box(slide, x, Inches(2.5), Inches(1.9), Inches(1.0),
                fill=BRAND_RED if i > 0 else ACCENT_YELLOW, line=None)
        add_text(slide, step, x, Inches(2.6), Inches(1.9), Inches(0.8),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # 箭头
        if i < len(flow_steps) - 1:
            add_text(slide, "→", x + Inches(1.9), Inches(2.7), Inches(0.3), Inches(0.5),
                     font_size=20, color=WHITE, align=PP_ALIGN.CENTER)

    # 杀手级应用
    add_text(slide, "杀手级应用：6-24 个月癌前病变预警", Inches(0.5), Inches(4.5), Inches(12), Inches(0.5),
             font_size=24, bold=True, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)

    kpi_data = [
        ("检出率", "30-40%", "70-80%"),
        ("预警提前", "0-6 月", "6-24 月"),
        ("误报率", "10-20%", "15-25%"),
    ]
    for i, (label, before, after) in enumerate(kpi_data):
        x = Inches(1 + i * 4)
        add_text(slide, label, x, Inches(5.2), Inches(3.5), Inches(0.4),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, f"{before} → {after}", x, Inches(5.7), Inches(3.5), Inches(0.6),
                 font_size=20, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)

    # ===== Slide 4: 理论基础 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_text(slide, "理论基础：张颖清全息胚理论（1981）", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=36, bold=True, color=BRAND_RED)

    # 引言
    quote = (
        "生物体每一相对独立的部分（耳、第二掌骨、足底等）\n"
        "都包含整体各部位的全部信息。"
    )
    add_text(slide, quote, Inches(1), Inches(1.5), Inches(11), Inches(1.0),
             font_size=20, italic=True, color=DARK_BG, align=PP_ALIGN.CENTER)

    # 三大支柱
    pillars = [
        ("临床验证", "中医 3000 年\n200+ 种全息胚疗法"),
        ("现代医学证实", "Dermatome / Head's Zones\n牵涉痛（神经解剖学）"),
        ("跨文化验证", "阿育吠陀 Marma / 藏医 / 韩医\nUnani（60-80% 重叠）"),
    ]
    for i, (title, desc) in enumerate(pillars):
        x = Inches(0.5 + i * 4.3)
        add_box(slide, x, Inches(3.5), Inches(4), Inches(3),
                fill=BRAND_RED, line=None)
        add_text(slide, title, x, Inches(3.8), Inches(4), Inches(0.6),
                 font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x, Inches(4.8), Inches(4), Inches(1.5),
                 font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

    # 底部
    add_text(slide, "核心论文：张颖清《生物全息诊疗法》（1985）",
             Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
             font_size=14, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # ===== Slide 5: 市场规模 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_text(slide, "市场规模", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    # TAM 饼图（用堆叠条形）
    markets = [
        ("中国中医现代化", 1000, ACCENT_YELLOW),
        ("全球自然疗法", 3000, ACCENT_BLUE),
        ("健康监测硬件", 5000, BRAND_RED),
        ("保险+药企数据", 500, WHITE),
    ]
    total = 9500
    y = 1.5
    for label, size, color in markets:
        pct = size / total * 100
        bar_width = Inches(pct * 0.12)
        add_box(slide, Inches(2.5), Inches(y), bar_width, Inches(0.6),
                fill=color, line=None)
        add_text(slide, f"{label}: ¥{size} 亿 ({pct:.1f}%)", Inches(2.5), Inches(y), Inches(11), Inches(0.6),
                 font_size=18, color=WHITE, align=PP_ALIGN.LEFT)
        y += 0.9

    # 总 TAM
    add_text(slide, "总 TAM（可达市场）: ¥9,500 亿",
             Inches(0.5), Inches(5.5), Inches(12), Inches(0.6),
             font_size=28, bold=True, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)
    add_text(slide, "SAM（可服务市场）: ¥1,000 亿（中国中医 + AI）",
             Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
             font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "SOM（3 年目标）: ¥100 亿（5% SAM）",
             Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
             font_size=18, color=WHITE, align=PP_ALIGN.CENTER)

    # ===== Slide 6: 产品演示 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_text(slide, "产品演示", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    add_text(slide, "[App 截图]", Inches(0.5), Inches(1.5), Inches(12), Inches(0.6),
             font_size=24, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # 4 个 mockup 位置
    add_box(slide, Inches(1), Inches(2.5), Inches(5), Inches(4),
            fill=LIGHT_BG, line=GRAY)
    add_text(slide, "主屏幕：8 微系统选择", Inches(1), Inches(2.5), Inches(5), Inches(4),
             font_size=18, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(7), Inches(2.5), Inches(5), Inches(4),
            fill=LIGHT_BG, line=GRAY)
    add_text(slide, "耳图像分析结果", Inches(7), Inches(2.5), Inches(5), Inches(4),
             font_size=18, color=DARK_BG, align=PP_ALIGN.CENTER)

    add_text(slide, "技术规格：iOS / Android · 端侧 AI · 隐私保护（联邦学习）",
             Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
             font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # ===== Slide 7: 8 微系统 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_text(slide, "8 微系统 · 全球 60+ 疗法精华集成", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=36, bold=True, color=BRAND_RED)

    # 8 微系统表格
    microsystems = [
        ("耳", "耳廓 200+ 穴位", "Nogier 1957"),
        ("掌", "第二掌骨节肢", "张颖清 1973"),
        ("足", "足底 70+ 反射区", "Ingham 1938"),
        ("面", "面部 5 区", "中医《灵枢》"),
        ("头", "头皮针/EEG", "焦顺发 1971"),
        ("舌", "舌 5 区", "中医几千年"),
        ("脉", "桡动脉 28 脉象", "中医几千年"),
        ("EEG", "5 频段", "现代神经科学"),
    ]
    for i, (mod, area, src) in enumerate(microsystems):
        row = i // 4
        col = i % 4
        x = Inches(0.5 + col * 3.2)
        y = Inches(1.5 + row * 2.5)
        add_box(slide, x, y, Inches(3), Inches(2.2),
                fill=BRAND_RED, line=None)
        add_text(slide, mod, x, y + Inches(0.1), Inches(3), Inches(0.6),
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, area, x, y + Inches(0.8), Inches(3), Inches(0.6),
                 font_size=14, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, src, x, y + Inches(1.4), Inches(3), Inches(0.6),
                 font_size=12, italic=True, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)

    # ===== Slide 8: 技术架构 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_text(slide, "技术架构", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    layers = [
        ("Layer 4 · AI 辨证层", "多模态融合 Transformer · 中医辨证算法 · 个性化基线模型"),
        ("Layer 3 · 数字孪生层", "用户健康状态虚拟副本 · 时间序列追踪 · 异常检测"),
        ("Layer 2 · AI 模型层", "8 微系统专用模型 · 端侧轻量化 · 联邦学习"),
        ("Layer 1 · 数据采集层", "手机摄像头 + 便携 EEG · 标准化光源 · 实时传输"),
    ]
    for i, (title, desc) in enumerate(layers):
        y = Inches(1.3 + i * 1.4)
        add_box(slide, Inches(1), y, Inches(11), Inches(1.1),
                fill=BRAND_RED if i % 2 == 0 else DARK_BG, line=WHITE)
        add_text(slide, title, Inches(1.2), y + Inches(0.1), Inches(10), Inches(0.4),
                 font_size=20, bold=True, color=ACCENT_YELLOW)
        add_text(slide, desc, Inches(1.2), y + Inches(0.6), Inches(10), Inches(0.4),
                 font_size=14, color=WHITE)

    # ===== Slide 9: 杀手级应用 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_text(slide, "杀手级应用", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)
    add_text(slide, "6-24 个月癌前病变预警", Inches(0.5), Inches(1.1), Inches(12), Inches(0.6),
             font_size=28, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)

    # 对比表
    rows = [
        ("早期检出率", "30-40%", "70-80%", True),
        ("预警提前量", "0-6 月", "6-24 月", True),
        ("假阳性率", "10-20%", "15-25%", False),
        ("个性化", "群体标准", "个体基线", True),
        ("覆盖", "单器官", "8 微系统", True),
    ]
    # 表头
    add_text(slide, "指标", Inches(1), Inches(2.3), Inches(3), Inches(0.5),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "当前临床", Inches(5), Inches(2.3), Inches(3), Inches(0.5),
             font_size=18, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "HoloScan 2.0", Inches(8.5), Inches(2.3), Inches(3.5), Inches(0.5),
             font_size=18, bold=True, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)

    for i, (label, before, after, highlight) in enumerate(rows):
        y = Inches(2.9 + i * 0.7)
        color = ACCENT_YELLOW if highlight else WHITE
        add_text(slide, label, Inches(1), y, Inches(3), Inches(0.5),
                 font_size=16, color=color)
        add_text(slide, before, Inches(5), y, Inches(3), Inches(0.5),
                 font_size=16, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, after, Inches(8.5), y, Inches(3.5), Inches(0.5),
                 font_size=18, bold=highlight, color=color, align=PP_ALIGN.CENTER)

    add_text(slide, "基于 6.5 突变检测（HoloScan 2.0 杀手级模块）",
             Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # ===== Slide 10: 商业模式 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_text(slide, "商业模式", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    # 表头
    headers = ["模式", "价格", "目标", "Y3 收入", "Y5 收入"]
    for i, h in enumerate(headers):
        x = Inches(0.5 + i * 2.4)
        add_text(slide, h, x, Inches(1.3), Inches(2.4), Inches(0.5),
                 font_size=16, bold=True, color=BRAND_RED, align=PP_ALIGN.CENTER)

    rows = [
        ("C 端订阅", "¥99/月", "个人健康", "¥600 万", "¥3,000 万"),
        ("B 端诊所", "¥5,000/月", "中医诊所", "¥1,500 万", "¥1.5 亿"),
        ("数据 API", "¥0.1/次", "保险公司", "¥300 万", "¥3,000 万"),
        ("企业定制", "¥10-100 万", "大型医院", "¥600 万", "¥6,000 万"),
    ]
    for i, row in enumerate(rows):
        y = Inches(1.9 + i * 0.8)
        bg = LIGHT_BG if i % 2 == 0 else WHITE
        for j, cell in enumerate(row):
            x = Inches(0.5 + j * 2.4)
            add_box(slide, x, y, Inches(2.4), Inches(0.7), fill=bg, line=GRAY)
            add_text(slide, cell, x, y, Inches(2.4), Inches(0.7),
                     font_size=14, color=DARK_BG, align=PP_ALIGN.CENTER)

    # 总计
    add_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.2),
            fill=BRAND_RED, line=None)
    add_text(slide, "总收入：Y3 ¥3,000 万  →  Y5 ¥2.3 亿",
             Inches(0.5), Inches(5.6), Inches(12), Inches(1),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ===== Slide 11: 当前进展 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_text(slide, "当前进展", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    achievements = [
        "✅ 6.x 物理学基础系列（~2 万字）",
        "✅ 7.x 应用扩展系列（~26.5K 字）",
        "✅ 8.x 项目计划 + agency-agents 协作框架",
        "✅ HoloAgent v0.2.6（生产化 AI Agent）",
        "✅ HAIS v0.1（AI Agent 接口标准）",
        "✅ 1000 页专著提纲",
        "✅ 30 集短视频课程完整脚本",
        "✅ HoloScan 商业构想 v0.1",
    ]
    for i, item in enumerate(achievements):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.3)
        y = Inches(1.5 + row * 0.8)
        add_text(slide, item, x, y, Inches(6), Inches(0.6),
                 font_size=18, color=DARK_BG, align=PP_ALIGN.LEFT)

    # 核心数据
    add_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5),
            fill=BRAND_RED, line=None)
    add_text(slide, "核心文档：~17 万字 / 44+ 文档",
             Inches(0.5), Inches(5.6), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "覆盖 6 大传统医学 / 200+ 疗法 / 8 微系统 / 1 万页研究内容",
             Inches(0.5), Inches(6.3), Inches(12), Inches(0.6),
             font_size=18, color=WHITE, align=PP_ALIGN.CENTER)

    # ===== Slide 12: 团队 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_text(slide, "团队", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    # 核心团队
    team = [
        ("创始人 / CEO", "用户本人", "项目主导者"),
        ("AI 工程师", "待聘", "多模态融合"),
        ("临床总监", "待聘", "中医 + 西医"),
        ("监管总监", "待聘", "NMPA / FDA"),
        ("产品总监", "待聘", "产品规划"),
        ("市场总监", "待聘", "营销"),
    ]
    for i, (role, name, desc) in enumerate(team):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.3)
        y = Inches(1.5 + row * 2.0)
        add_box(slide, x, y, Inches(4), Inches(1.8),
                fill=BRAND_RED, line=None)
        add_text(slide, role, x, y + Inches(0.1), Inches(4), Inches(0.5),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, name, x, y + Inches(0.7), Inches(4), Inches(0.4),
                 font_size=14, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x, y + Inches(1.2), Inches(4), Inches(0.4),
                 font_size=12, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    # AI Agent 团队
    add_text(slide, "+ AI Agent 团队：31 个 Agent 跨 8 部门",
             Inches(0.5), Inches(5.8), Inches(12), Inches(0.6),
             font_size=20, bold=True, color=BRAND_RED, align=PP_ALIGN.CENTER)
    add_text(slide, "工程 / 设计 / 产品 / 市场 / 项目管理 / 测试 / 支持 / 特别",
             Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # ===== Slide 13: 路线图 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_text(slide, "路线图", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    milestones = [
        ("2026 Q3-Q4", "立项", "2 月", "M1 架构完成"),
        ("2026 Q4-2027 Q1", "MVP", "6 月", "M2 耳模块 + 100 用户"),
        ("2027 Q1-2028 Q1", "完整版", "12 月", "M3 8 微系统 + 1000 用户"),
        ("2028 Q1-2029 Q1", "临床", "12 月", "M4-M5 6 项 RCT"),
        ("2029 Q1-Q3", "商业化", "6 月", "M6-M8 上市 + ¥3,000 万"),
    ]
    for i, (period, name, duration, deliverable) in enumerate(milestones):
        y = Inches(1.3 + i * 1.0)
        # 时间标签
        add_box(slide, Inches(0.5), y, Inches(2.5), Inches(0.8),
                fill=BRAND_RED, line=None)
        add_text(slide, period, Inches(0.5), y, Inches(2.5), Inches(0.8),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # 阶段名
        add_text(slide, name, Inches(3.2), y, Inches(1.5), Inches(0.8),
                 font_size=20, bold=True, color=ACCENT_YELLOW)
        # 时长
        add_text(slide, duration, Inches(4.8), y, Inches(1), Inches(0.8),
                 font_size=14, color=WHITE, align=PP_ALIGN.CENTER)
        # 交付
        add_text(slide, deliverable, Inches(6), y, Inches(6.5), Inches(0.8),
                 font_size=14, color=WHITE)

    # ===== Slide 14: 竞争优势 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_text(slide, "竞争优势 · 5 大壁垒", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=36, bold=True, color=BRAND_RED)

    advantages = [
        ("理论先发", "基于张颖清 1985 理论 + 现代神经科学"),
        ("多模态壁垒", "8 微系统集成需要 ~5 年研发"),
        ("临床数据壁垒", "癌前病变前瞻队列 5000 例"),
        ("学术壁垒", "Nature/Lancet 论文 + 全球全息胚数据库"),
        ("跨学科壁垒", "中医 + 神经 + AI + 工程团队"),
    ]
    for i, (title, desc) in enumerate(advantages):
        y = Inches(1.5 + i * 0.9)
        add_text(slide, f"#{i+1} {title}", Inches(0.5), y, Inches(3), Inches(0.7),
                 font_size=20, bold=True, color=BRAND_RED)
        add_text(slide, desc, Inches(3.8), y, Inches(9), Inches(0.7),
                 font_size=16, color=DARK_BG, align=PP_ALIGN.LEFT)

    # vs 竞争对手
    add_box(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(1.3),
            fill=LIGHT_BG, line=GRAY)
    add_text(slide, "vs 竞争对手",
             Inches(0.5), Inches(6.0), Inches(12), Inches(0.4),
             font_size=18, bold=True, color=BRAND_RED, align=PP_ALIGN.CENTER)
    add_text(slide, "• vs 中医 AI 诊断（单一证型）→ 多模态融合优势",
             Inches(0.5), Inches(6.4), Inches(12), Inches(0.3),
             font_size=14, color=DARK_BG)
    add_text(slide, "• vs Reflexology App（仅足底）→ 8 微系统优势",
             Inches(0.5), Inches(6.7), Inches(12), Inches(0.3),
             font_size=14, color=DARK_BG)
    add_text(slide, "• vs 穿戴设备（仅生理）→ 中医辨证优势",
             Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
             font_size=14, color=DARK_BG)

    # ===== Slide 15: 财务预测 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_text(slide, "财务预测", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    # 5 年财务
    headers = ["年", "用户", "收入", "关键"]
    rows = [
        ("Y1", "0", "¥0", "开发期"),
        ("Y2", "100", "¥500 万", "早期试用"),
        ("Y3", "1,000", "¥3,000 万", "MVP + 上市"),
        ("Y4", "5,000", "¥1 亿", "临床验证完成"),
        ("Y5", "20,000", "¥3 亿", "完整商业化"),
    ]
    # 表头
    add_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
            fill=BRAND_RED, line=None)
    for i, h in enumerate(headers):
        x = Inches(1.5 + i * 2.5)
        add_text(slide, h, x, Inches(1.5), Inches(2.5), Inches(0.6),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        y = Inches(2.3 + i * 0.7)
        bg = LIGHT_BG if i % 2 == 0 else WHITE
        add_box(slide, Inches(1), y, Inches(11), Inches(0.6),
                fill=bg, line=GRAY)
        for j, cell in enumerate(row):
            x = Inches(1.5 + j * 2.5)
            color = ACCENT_YELLOW if j == 2 and i >= 2 else DARK_BG
            add_text(slide, cell, x, y, Inches(2.5), Inches(0.6),
                     font_size=18, bold=(j == 2), color=color, align=PP_ALIGN.CENTER)

    # 关键指标
    add_text(slide, "毛利率：75-85% | LTV/CAC：12-24x（健康）",
             Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
             font_size=20, bold=True, color=BRAND_RED, align=PP_ALIGN.CENTER)

    # ===== Slide 16: 融资需求 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_text(slide, "融资需求", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    add_text(slide, "种子轮：¥500 万",
             Inches(0.5), Inches(1.5), Inches(12), Inches(1.2),
             font_size=72, bold=True, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)
    add_text(slide, "估值 ¥5,000 万（pre-money）",
             Inches(0.5), Inches(2.8), Inches(12), Inches(0.6),
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    # 资金用途
    usage = [
        ("团队", "¥300 万", 60),
        ("服务器/AI 训练", "¥100 万", 20),
        ("用户测试", "¥50 万", 10),
        ("注册/法务", "¥30 万", 6),
        ("杂项", "¥20 万", 4),
    ]
    add_text(slide, "资金用途", Inches(0.5), Inches(4.0), Inches(12), Inches(0.5),
             font_size=20, bold=True, color=WHITE)
    for i, (label, amount, pct) in enumerate(usage):
        y = Inches(4.7 + i * 0.4)
        bar_w = Inches(pct * 0.12)
        add_box(slide, Inches(0.5), y, bar_w, Inches(0.3), fill=ACCENT_YELLOW, line=None)
        add_text(slide, f"{label}: {amount} ({pct}%)", Inches(0.5), y, Inches(11), Inches(0.3),
                 font_size=14, color=WHITE, align=PP_ALIGN.LEFT)

    add_text(slide, "里程碑：M2 MVP 发布 + 100 用户（NPS 50+）",
             Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
             font_size=16, italic=True, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)

    # ===== Slide 17: 资金使用 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_text(slide, "资金使用分配", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    # 饼图（用方块表示）
    allocations = [
        ("团队 60%", "¥300 万", BRAND_RED),
        ("AI 算力 20%", "¥100 万", ACCENT_YELLOW),
        ("用户测试 10%", "¥50 万", ACCENT_BLUE),
        ("注册法务 6%", "¥30 万", GRAY),
        ("杂项 4%", "¥20 万", WHITE),
    ]
    for i, (label, amount, color) in enumerate(allocations):
        row = i // 3
        col = i % 3
        x = Inches(0.5 + col * 4.3)
        y = Inches(1.5 + row * 2.5)
        # 圆角矩形（用 box 替代）
        add_box(slide, x, y, Inches(4), Inches(2.2), fill=color, line=DARK_BG)
        add_text(slide, label, x, y + Inches(0.3), Inches(4), Inches(0.7),
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, amount, x, y + Inches(1.0), Inches(4), Inches(0.7),
                 font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ===== Slide 18: 退出策略 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_text(slide, "退出策略", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    exit_paths = [
        ("A 轮融资", "M24", "¥3 亿", "6x"),
        ("B 轮融资", "M36", "¥10 亿", "20x"),
        ("Pre-IPO", "Y5", "¥30 亿", "60x"),
        ("战略并购", "Y5+", "¥50-100 亿", "100-200x"),
        ("IPO", "Y7+", "¥100-300 亿", "200-600x"),
    ]
    # 表头
    add_box(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(0.6),
            fill=BRAND_RED, line=None)
    for i, h in enumerate(["路径", "时点", "估值", "倍数"]):
        x = Inches(0.5 + i * 3)
        add_text(slide, h, x, Inches(1.5), Inches(3), Inches(0.6),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, row in enumerate(exit_paths):
        y = Inches(2.3 + i * 0.6)
        bg = LIGHT_BG if i % 2 == 0 else WHITE
        add_box(slide, Inches(0.5), y, Inches(12), Inches(0.5),
                fill=bg, line=GRAY)
        for j, cell in enumerate(row):
            x = Inches(0.5 + j * 3)
            add_text(slide, cell, x, y, Inches(3), Inches(0.5),
                     font_size=14, color=DARK_BG, align=PP_ALIGN.CENTER)

    add_text(slide, "潜在并购方：平安好医生 / 罗氏 / 阿里健康",
             Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
             font_size=16, italic=True, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)

    # ===== Slide 19: 愿景 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_text(slide, "愿景", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=44, bold=True, color=BRAND_RED)

    add_text(slide, "让全息胚理论成为 21 世纪统一医学的工程基础",
             Inches(0.5), Inches(1.5), Inches(12), Inches(1),
             font_size=32, bold=True, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)

    # 短中长期
    timeline = [
        ("短期（5 年）", "全球最大全息胚 AI 平台", "100 万付费用户", "年收入 ¥3-10 亿"),
        ("中期（10 年）", "全球统一医学标杆", "1000 万用户", "年收入 ¥30-100 亿"),
        ("长期（20 年）", "AGI 时代的全息胚系统", "太空医学标配", "重新定义'健康'"),
    ]
    for i, (period, target, users, income) in enumerate(timeline):
        y = Inches(3.0 + i * 1.3)
        add_box(slide, Inches(0.5), y, Inches(12), Inches(1.1),
                fill=BRAND_RED, line=None)
        add_text(slide, period, Inches(0.5), y, Inches(2.5), Inches(1.1),
                 font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, target, Inches(3.0), y + Inches(0.1), Inches(4), Inches(0.5),
                 font_size=16, bold=True, color=WHITE)
        add_text(slide, users, Inches(3.0), y + Inches(0.55), Inches(4), Inches(0.5),
                 font_size=14, color=WHITE)
        add_text(slide, income, Inches(7.0), y + Inches(0.3), Inches(5.5), Inches(0.5),
                 font_size=16, bold=True, color=ACCENT_YELLOW)

    # ===== Slide 20: 联系我们 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_text(slide, "HoloScan 2.0", Inches(1), Inches(2), Inches(11), Inches(1.5),
             font_size=72, bold=True, color=BRAND_RED, align=PP_ALIGN.CENTER)
    add_text(slide, "让每一次体检都更早一步", Inches(1), Inches(3.5), Inches(11), Inches(0.8),
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "创始人：[用户姓名]", Inches(1), Inches(5.0), Inches(11), Inches(0.5),
             font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "邮箱：[待补充] · 电话：[待补充]", Inches(1), Inches(5.6), Inches(11), Inches(0.5),
             font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "地址：厦门", Inches(1), Inches(6.2), Inches(11), Inches(0.5),
             font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "项目地址：D:\\obsidian\\Holographic-Biology\\", Inches(1), Inches(6.8), Inches(11), Inches(0.4),
             font_size=14, italic=True, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)

    # 保存
    prs.save(OUTPUT)
    print(f"✅ 已生成 PPTX：{OUTPUT}")
    print(f"   共 {len(prs.slides)} 页")


# 辅助函数
def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=None, align=None, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color


def add_box(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


if __name__ == "__main__":
    create_presentation()